#!/usr/bin/env python
# encoding: utf-8

# Copyright (C) 2016 Steven Myint
#
# Permission is hereby granted, free of charge, to any person obtaining
# a copy of this software and associated documentation files (the
# "Software"), to deal in the Software without restriction, including
# without limitation the rights to use, copy, modify, merge, publish,
# distribute, sublicense, and/or sell copies of the Software, and to
# permit persons to whom the Software is furnished to do so, subject to
# the following conditions:
#
# The above copyright notice and this permission notice shall be
# included in all copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND,
# EXPRESS OR IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF
# MERCHANTABILITY, FITNESS FOR A PARTICULAR PURPOSE AND
# NONINFRINGEMENT. IN NO EVENT SHALL THE AUTHORS OR COPYRIGHT HOLDERS
# BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER LIABILITY, WHETHER IN AN
# ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM, OUT OF OR IN
# CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

"""Converts reStructuredText to PowerPoint."""

import io
import os
import sys
import urllib

import docutils.core
import docutils.nodes
import docutils.utils
import pptx


__version__ = '0.3'


TITLE_BUFFER = pptx.util.Inches(2.)
MARGIN = pptx.util.Inches(1.)


class PowerPointTranslator(docutils.nodes.NodeVisitor):

    """A translator for converting docutils elements to PowerPoint."""

    def __init__(self, document, presentation):
        docutils.nodes.NodeVisitor.__init__(self, document)

        self.bullet_level = 0
        self.presentation = presentation
        self.slides = self.presentation.slides
        self.table_rows = None
        self.title_slide = True

    def visit_document(self, node):
        pass

    def depart_document(self, node):
        pass

    def visit_docinfo_item(self, node, name):
        pass

    def visit_image(self, node):
        uri = node.attributes['uri']
        if '://' in uri:
            if sys.version_info[0] < 3:
                self.document.reporter.warning(
                    'Downloading images requires Python 3 or greater')
                return

            try:
                with urllib.request.urlopen(uri) as input_file:
                    image_file = io.BytesIO(input_file.read())
            except urllib.error.HTTPError as e:
                self.document.reporter.warning(
                    'Could not access {}'.format(uri))
                return
        else:
            document_filename = docutils.utils.get_source_line(node)[0]
            if document_filename and document_filename != '<stdin>':
                root_path = os.path.dirname(document_filename)
            else:
                root_path = os.getcwd()
            image_file = os.path.join(root_path, uri)
            uri = image_file

        try:
            picture = self.slides[-1].shapes.add_picture(
                image_file,
                left=0,
                top=0)
        except IOError:
            self.document.reporter.warning(
                'Could not load image {}'.format(uri))
            return

        center_picture(picture, self.presentation)

    def depart_image(self, node):
        pass

    def visit_Text(self, node):
        pass

    def depart_Text(self, node):
        pass

    def visit_list_item(self, node):
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.add_paragraph()
        paragraph.text = node.astext()

        if self.bullet_level:
            paragraph.level = self.bullet_level

        raise docutils.nodes.SkipNode

    def depart_list_item(self, node):
        pass

    def visit_paragraph(self, node):
        shapes = self.slides[-1].shapes

        if self.title_slide and not shapes[-1].text:
            # This must be the empty text box for the subtitle.
            shapes[-1].text = node.astext()
        else:
            text_frame = self.slides[-1].shapes.placeholders[1].text_frame
            paragraph = text_frame.add_paragraph()
            paragraph.text = node.astext()

    def depart_paragraph(self, node):
        pass

    def visit_section(self, node):
        self.title_slide = False
        self.slides.add_slide(self.presentation.slide_layouts[1])

    def depart_section(self, node):
        pass

    def visit_title(self, node):
        if len(self.slides):
            self.slides[-1].shapes.title.text = node.astext()
        else:
            # Title slide.
            slide = self.slides.add_slide(self.presentation.slide_layouts[0])
            slide.shapes.title.text = node.astext()
            self.title_slide = True
            # TODO: Author.

    def depart_title(self, node):
        pass

    def visit_literal_block(self, node):
        pass

    def depart_literal_block(self, node):
        pass

    def visit_bullet_list(self, node):
        self.bullet_level += 1

    def depart_bullet_list(self, node):
        self.bullet_level -= 1
        assert self.bullet_level >= 0

    def visit_enumerated_list(self, node):
        pass

    def depart_enumerated_list(self, node):
        pass

    def visit_tgroup(self, node):
        self.table_rows = []

    def depart_tgroup(self, node):
        if self.table_rows and self.table_rows[0]:
            table = self.slides[-1].shapes.add_table(
                rows=len(self.table_rows),
                cols=len(self.table_rows[0]),
                left=MARGIN,
                top=TITLE_BUFFER,
                width=self.presentation.slide_width - 2 * MARGIN,
                height=self.presentation.slide_height - 2 * TITLE_BUFFER).table

            for (row_index, row) in enumerate(self.table_rows):
                for (col_index, col) in enumerate(row):
                    table.cell(row_idx=row_index, col_idx=col_index).text = col

            self.table_rows = None

    def visit_row(self, node):
        assert self.table_rows is not None
        self.table_rows.append([])

    def depart_row(self, node):
        pass

    def visit_entry(self, node):
        self.table_rows[-1].append(node.astext())
        raise docutils.nodes.SkipNode

    def depart_entry(self, node):
        pass

    def unknown_visit(self, node):
        self.document.reporter.warning('unknown_visit({})'.format(node))

    def unknown_departure(self, node):
        self.document.reporter.warning('unknown_departure({})'.format(node))

    def astext(self):
        pass


def center_picture(picture, presentation):
    picture.left = (presentation.slide_width - picture.width) // 2

    slide_height = presentation.slide_height - TITLE_BUFFER
    picture.top = (slide_height - picture.height) // 2 + TITLE_BUFFER


class PowerPointWriter(docutils.core.writers.Writer):

    """A docutils writer that produces PowerPoint."""

    settings_spec = (
        'PowerPoint options',
        None,
        (
            (
                'PowerPoint template.',
                ['--pptx-template'],
                {'default': None}
            ),
        )
    )

    def __init__(self):
        docutils.core.writers.Writer.__init__(self)

        self.presentation = None
        self.translator_class = PowerPointTranslator

    def translate(self):
        assert self.presentation
        visitor = self.translator_class(document=self.document,
                                        presentation=self.presentation)
        self.document.walkabout(visitor)

    def write(self, document, destination):
        self.document = document
        self.presentation = pptx.Presentation(document.settings.pptx_template)

        self.language = docutils.languages.get_language(
            document.settings.language_code,
            document.reporter)

        self.translate()

        if destination.destination is None:
            self.presentation.save(destination.destination_path)
        else:
            stream = io.BytesIO()
            self.presentation.save(stream)
            destination.write(stream.getvalue())


def main():
    docutils.core.publish_cmdline(
        writer=PowerPointWriter(),
        description='Generates PowerPoint presentations.  ' +
                    docutils.core.default_description,
        settings_overrides={'halt_level': docutils.utils.Reporter.ERROR_LEVEL})


if __name__ == '__main__':
    main()
